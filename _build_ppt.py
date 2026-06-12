"""
연결재무보고 통합 시스템 — 임원 보고용 PPT 생성기 (v2: 31장)
- python-pptx 사용
- 16:9 와이드스크린, 31장 (신규 6개 화면 반영)
- 화면 캡쳐는 'ppt_screenshots' 폴더에 저장된 파일을 자동 인식하여 삽입
  (없으면 회색 placeholder 박스 + 안내 문구 표시)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(BASE_DIR, "연결재무보고시스템_임원보고용.pptx")
SHOTS_DIR = os.path.join(BASE_DIR, "ppt_screenshots")
os.makedirs(SHOTS_DIR, exist_ok=True)

# ===== 색상 팔레트 (Midnight Executive) =====
NAVY      = RGBColor(0x1E, 0x27, 0x61)
NAVY_DARK = RGBColor(0x14, 0x1B, 0x47)
ICE       = RGBColor(0xCA, 0xDC, 0xFC)
ACCENT    = RGBColor(0xF9, 0xB7, 0x1A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DK   = RGBColor(0x33, 0x33, 0x33)
GRAY_MID  = RGBColor(0x6B, 0x72, 0x80)
GRAY_LT   = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_BG   = RGBColor(0xF5, 0xF6, 0xFA)

FONT_HEAD = "맑은 고딕"
FONT_BODY = "맑은 고딕"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]

# ===== 헬퍼 =====
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=GRAY_DK,
             font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=14, color=GRAY_DK,
                font=FONT_BODY, bullet_color=NAVY, line_spacing=1.35):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            txt, opts = item
        else:
            txt, opts = item, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = "■  "
        r1.font.name = font
        r1.font.size = Pt(opts.get("size", size))
        r1.font.color.rgb = opts.get("bullet_color", bullet_color)
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = txt
        r2.font.name = font
        r2.font.size = Pt(opts.get("size", size))
        r2.font.color.rgb = opts.get("color", color)
        r2.font.bold = opts.get("bold", False)
    return tb

def add_title_bar(slide, title, subtitle=None, num=None, total=None):
    add_rect(slide, Inches(0.5), Inches(0.55), Inches(0.12), Inches(0.55), NAVY)
    add_text(slide, Inches(0.75), Inches(0.45), Inches(10.5), Inches(0.55),
             title, size=26, bold=True, color=NAVY_DARK, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.75), Inches(1.02), Inches(10.5), Inches(0.32),
                 subtitle, size=13, color=GRAY_MID)
    if num is not None:
        add_text(slide, Inches(11.5), Inches(0.5), Inches(1.4), Inches(0.4),
                 f"{num:02d} / {total:02d}", size=10, color=GRAY_MID,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(0.5), Inches(7.1), Inches(12.33), Emu(9525), GRAY_LT)
    add_text(slide, Inches(0.5), Inches(7.18), Inches(8), Inches(0.25),
             "연결재무보고 통합 시스템 | 임원 보고용", size=9, color=GRAY_MID)

def find_screenshot(keys):
    if not os.path.isdir(SHOTS_DIR):
        return None
    files = os.listdir(SHOTS_DIR)
    exts = ('.png', '.jpg', '.jpeg', '.gif', '.bmp')
    for k in keys:
        for f in files:
            if f.lower().endswith(exts) and (f.lower().startswith(k.lower()) or k.lower() in f.lower()):
                return os.path.join(SHOTS_DIR, f)
    return None

def add_screenshot_area(slide, x, y, w, h, slide_num, label, keys=None):
    keys = keys or []
    keys = [f"{slide_num:02d}"] + keys
    path = find_screenshot(keys)
    add_rect(slide, x, y, w, h, GRAY_BG, line=GRAY_LT)
    if path and os.path.exists(path):
        try:
            from PIL import Image
            with Image.open(path) as im:
                iw, ih = im.size
            ratio_img  = iw / ih
            ratio_box  = w / h
            if ratio_img > ratio_box:
                nw = w
                nh = int(w / ratio_img)
                nx = x
                ny = y + (h - nh) // 2
            else:
                nh = h
                nw = int(h * ratio_img)
                ny = y
                nx = x + (w - nw) // 2
            slide.shapes.add_picture(path, nx, ny, width=nw, height=nh)
            add_text(slide, x, y + h + Inches(0.05), w, Inches(0.25),
                     f"〈 {label} 〉", size=10, color=GRAY_MID, align=PP_ALIGN.CENTER, italic=True)
        except Exception as e:
            add_text(slide, x, y, w, h,
                     f"[이미지 로드 실패]\n{path}\n{e}", size=12, color=RGBColor(0xC0, 0x39, 0x2B),
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:
        add_text(slide, x, y, w, h - Inches(0.6),
                 "[ 화면 캡쳐 삽입 ]", size=24, bold=True, color=RGBColor(0xB0, 0xB7, 0xC3),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        guide = f'파일명 예시:  {slide_num:02d}_{(keys[1] if len(keys)>1 else "screenshot")}.png\n저장 경로:  ppt_screenshots\\'
        add_text(slide, x, y + h - Inches(0.7), w, Inches(0.6),
                 guide, size=10, color=GRAY_MID, align=PP_ALIGN.CENTER)
        add_text(slide, x, y + h + Inches(0.05), w, Inches(0.25),
                 f"〈 {label} 〉", size=10, color=GRAY_MID, align=PP_ALIGN.CENTER, italic=True)

def add_tag(slide, x, y, label, color=NAVY, text_color=WHITE):
    w, h = Inches(0.7), Inches(0.28)
    add_rect(slide, x, y, w, h, color)
    add_text(slide, x, y, w, h, label, size=10, bold=True, color=text_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

TOTAL = 31

# ============================================================
# 슬라이드 1 : 표지
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SW, SH, NAVY_DARK)
add_rect(s, Inches(9.5), 0, Inches(3.83), SH, NAVY)
add_rect(s, Inches(9.3), 0, Inches(0.06), SH, ACCENT)
add_text(s, Inches(0.9), Inches(1.3), Inches(8), Inches(0.45),
         "EXECUTIVE BRIEFING", size=14, bold=True, color=ACCENT, font=FONT_HEAD)
add_text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(1.4),
         "연결재무보고 통합 시스템", size=54, bold=True, color=WHITE, font=FONT_HEAD)
add_text(s, Inches(0.9), Inches(3.2), Inches(11), Inches(0.7),
         "주요 기능 소개 (확장판)", size=30, color=ICE, font=FONT_HEAD)
add_rect(s, Inches(0.9), Inches(4.15), Inches(0.8), Emu(38100), ACCENT)
add_text(s, Inches(0.9), Inches(4.35), Inches(11), Inches(0.6),
         "30개+ 자회사 결산패키지 통합 · 자동 합산 · 연결조정 · 현금정산표 · 검증 자동화", size=15, color=ICE)
add_text(s, Inches(0.9), Inches(6.55), Inches(6), Inches(0.3),
         "재무팀 자체 개발 | 2026", size=11, color=ICE)
add_text(s, Inches(6.9), Inches(6.55), Inches(5.5), Inches(0.3),
         "임원 보고용", size=11, color=ICE, align=PP_ALIGN.RIGHT)

# ============================================================
# 슬라이드 2 : 목차
# ============================================================
s = prs.slides.add_slide(blank)
add_title_bar(s, "목차", "Contents", 2, TOTAL)

contents = [
    ("01", "도입 배경",         "현행 업무의 비효율과 리스크"),
    ("02", "시스템 개요",       "5단계 결산 프로세스 통합"),
    ("03", "주요 기능 소개",    "메인 · 권한 · 패키지 · WCE · 연결 · 검증"),
    ("04", "기대 효과",         "정량 · 정성 효과"),
    ("05", "기술 및 운영",      "보안 · 사양 · 배포 일정"),
    ("06", "도입 로드맵 및 협조 요청", "단계별 일정과 IT부서 협조 사항"),
]
card_w, card_h = Inches(5.9), Inches(1.55)
gap_x, gap_y = Inches(0.35), Inches(0.25)
start_x, start_y = Inches(0.6), Inches(1.7)
for i, (num, title, desc) in enumerate(contents):
    col, row = i % 2, i // 2
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_rect(s, x, y, card_w, card_h, WHITE, line=GRAY_LT)
    add_rect(s, x, y, Inches(0.1), card_h, NAVY)
    add_text(s, x + Inches(0.35), y + Inches(0.15), Inches(1.2), Inches(0.6),
             num, size=36, bold=True, color=NAVY, font=FONT_HEAD)
    add_text(s, x + Inches(1.55), y + Inches(0.25), card_w - Inches(1.7), Inches(0.5),
             title, size=18, bold=True, color=GRAY_DK)
    add_text(s, x + Inches(1.55), y + Inches(0.78), card_w - Inches(1.7), Inches(0.6),
             desc, size=12, color=GRAY_MID)

# ============================================================
# 슬라이드 3 : Executive Summary
# ============================================================
s = prs.slides.add_slide(blank)
add_title_bar(s, "Executive Summary", "1페이지 요약", 3, TOTAL)

add_rect(s, Inches(0.5), Inches(1.6), Inches(5.5), Inches(5.25), NAVY)
add_text(s, Inches(0.85), Inches(1.85), Inches(4.8), Inches(0.4),
         "KEY MESSAGE", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(2.35), Inches(4.8), Inches(2.5),
         ["30개+ 자회사", "연결결산 전 과정을", "한 화면에서 처리합니다."],
         size=26, bold=True, color=WHITE, line_spacing=1.25)
add_rect(s, Inches(0.85), Inches(4.95), Inches(0.6), Emu(28575), ACCENT)
add_text(s, Inches(0.85), Inches(5.15), Inches(4.8), Inches(1.6),
         ["수기 작업 · 휴먼 에러 · 속인화(屬人化)의",
          "구조적 한계를 시스템으로 해소하고,",
          "결산 신뢰성과 대응 속도를 동시에 끌어올립니다."],
         size=13, color=ICE, line_spacing=1.45)

right_x = Inches(6.3)
items = [
    ("01", "업무 효율",   "분기당 수일 단위의 수기 합산 공수 제거. 즉시 결과 산출.", "0D9488"),
    ("02", "결산 신뢰성", "분개 검증 · 균형 체크 · 코드 감사 자동화로 휴먼 에러 차단.", "1E2761"),
    ("03", "감사 대응",   "권한 분리, 이력 관리, 표준화된 절차로 통제 강화.",        "B45309"),
]
for i, (n, title, desc, hex_c) in enumerate(items):
    y = Inches(1.6 + i*1.78)
    c = RGBColor.from_string(hex_c)
    add_rect(s, right_x, y, Inches(6.5), Inches(1.6), WHITE, line=GRAY_LT)
    add_rect(s, right_x, y, Inches(0.1), Inches(1.6), c)
    add_text(s, right_x + Inches(0.4), y + Inches(0.2), Inches(0.9), Inches(0.5),
             n, size=22, bold=True, color=c, font=FONT_HEAD)
    add_text(s, right_x + Inches(1.3), y + Inches(0.25), Inches(4.8), Inches(0.45),
             title, size=17, bold=True, color=GRAY_DK)
    add_text(s, right_x + Inches(1.3), y + Inches(0.78), Inches(5.1), Inches(0.8),
             desc, size=12, color=GRAY_MID, line_spacing=1.4)

# ============================================================
# 슬라이드 4 : 현행 업무의 비효율
# ============================================================
s = prs.slides.add_slide(blank)
add_title_bar(s, "현행 업무의 비효율", "왜 시스템화가 필요한가 (1/2)", 4, TOTAL)

add_rect(s, Inches(0.5), Inches(1.6), Inches(4), Inches(5.3), GRAY_BG)
add_text(s, Inches(0.5), Inches(2.2), Inches(4), Inches(1.4),
         "30+", size=110, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font=FONT_HEAD)
add_text(s, Inches(0.5), Inches(3.85), Inches(4), Inches(0.5),
         "결산 대상 자회사 수", size=15, bold=True, color=GRAY_DK, align=PP_ALIGN.CENTER)
add_rect(s, Inches(2.0), Inches(4.45), Inches(1.0), Emu(28575), ACCENT)
add_text(s, Inches(0.5), Inches(4.65), Inches(4), Inches(1.8),
         ["분기마다 반복되는",
          "엑셀 수기 합산 작업",
          "분기당 수일~수주 소요"],
         size=13, color=GRAY_MID, align=PP_ALIGN.CENTER, line_spacing=1.6)

right_x = Inches(4.9)
inefficiencies = [
    ("반복 공수", "동일한 양식의 결산패키지를 분기마다 수십 회 수기 취합 · 검증"),
    ("표준화 미비", "자회사별 양식 · 계산식 차이로 검토 · 정정 작업 반복 발생"),
    ("환율 적용 누락", "통화 환산 누락 · 환율 변경 시 전체 재작업 위험"),
    ("결과 산출 지연", "합산 → 검토 → 정정 사이클이 길어 의사결정 적시성 저하"),
]
add_text(s, right_x, Inches(1.55), Inches(8), Inches(0.5),
         "주요 비효율 요소", size=18, bold=True, color=NAVY_DARK)
for i, (t, d) in enumerate(inefficiencies):
    y = Inches(2.1 + i*1.18)
    add_rect(s, right_x, y, Inches(8.0), Inches(1.0), WHITE, line=GRAY_LT)
    add_rect(s, right_x, y, Inches(0.08), Inches(1.0), NAVY)
    add_text(s, right_x + Inches(0.3), y + Inches(0.15), Inches(7.5), Inches(0.4),
             t, size=15, bold=True, color=NAVY_DARK)
    add_text(s, right_x + Inches(0.3), y + Inches(0.55), Inches(7.5), Inches(0.45),
             d, size=12, color=GRAY_MID)

# ============================================================
# 슬라이드 5 : 업무 리스크
# ============================================================
s = prs.slides.add_slide(blank)
add_title_bar(s, "업무 리스크", "왜 시스템화가 필요한가 (2/2)", 5, TOTAL)

risks = [
    ("01", "휴먼 에러",     "수작업 합산 · 복사붙여넣기에서 발생하는 금액 · 코드 오류",   "EF4444"),
    ("02", "차변·대변 불균형", "연결조정 분개의 균형 검증이 누락될 위험",                "F97316"),
    ("03", "다단계 누락",   "글로벌세아 ⊃ 상역 ⊃ 태림/GIT 등 재귀 합산의 누락 · 중복", "EAB308"),
    ("04", "노하우 속인화", "특정 담당자에게 결산 노하우가 집중되어 리스크 가중",         "8B5CF6"),
    ("05", "감사 대응 취약", "절차의 비표준화로 내·외부 감사 대응 비용 증가",            "0EA5E9"),
    ("06", "정정 비용",     "결산 후 분개 정정 · 재합산으로 인한 시간 · 인력 낭비",       "10B981"),
]
cw, ch = Inches(4.05), Inches(2.4)
gx, gy = Inches(0.15), Inches(0.18)
sx, sy = Inches(0.5), Inches(1.55)
for i, (n, t, d, hex_c) in enumerate(risks):
    col, row = i % 3, i // 3
    x = sx + col * (cw + gx)
    y = sy + row * (ch + gy)
    c = RGBColor.from_string(hex_c)
    add_rect(s, x, y, cw, ch, WHITE, line=GRAY_LT)
    add_rect(s, x, y, cw, Inches(0.5), c)
    add_text(s, x + Inches(0.3), y, Inches(0.6), Inches(0.5),
             n, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.0), y, cw - Inches(1.1), Inches(0.5),
             t, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.25), y + Inches(0.75), cw - Inches(0.5), Inches(1.5),
             d, size=12, color=GRAY_DK, line_spacing=1.45)

# ============================================================
# 슬라이드 6 : 시스템 개요 (5단계 프로세스)
# ============================================================
s = prs.slides.add_slide(blank)
add_title_bar(s, "시스템 개요", "5단계 결산 프로세스를 하나의 시스템으로 통합", 6, TOTAL)

steps = [
    ("01", "패키지 배포·업로드", "분기 빈 패키지 자동 배포\n자회사 결산자료 일괄 업로드"),
    ("02", "자동 합산",          "BS · PL · 현금흐름\n자동 합산 + 환율 적용"),
    ("03", "WCE / 분개 입력",    "내부거래 제거 · 연결조정\n분개 입력 및 검증"),
    ("04", "다단계 그룹 합산",   "재귀 합산 · 현금정산표\n누락 · 중복 위험 제거"),
    ("05", "재무제표 · 대시보드", "연결재무제표 산출\nKPI 차트 자동 시각화"),
]
n = len(steps)
total_w = Inches(12.33)
step_w = (total_w - Inches(0.4)) / n
sy = Inches(2.7)
sx0 = Inches(0.5)
add_rect(s, sx0 + Inches(0.2), sy + Inches(0.55), total_w - Inches(0.4), Emu(28575), ICE)

for i, (num, title, desc) in enumerate(steps):
    x = sx0 + i * step_w
    cx = x + step_w/2 - Inches(0.55)
    cy = sy
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(1.1), Inches(1.1))
    shp.fill.solid(); shp.fill.fore_color.rgb = NAVY
    shp.line.color.rgb = NAVY
    tf = shp.text_frame
    tf.margin_left=Emu(0); tf.margin_right=Emu(0); tf.margin_top=Emu(0); tf.margin_bottom=Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = FONT_HEAD; r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE
    add_text(s, x, sy + Inches(1.25), step_w, Inches(0.5),
             title, size=14, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.05), sy + Inches(1.75), step_w - Inches(0.1), Inches(1.6),
             desc.split("\n"), size=11, color=GRAY_MID, align=PP_ALIGN.CENTER, line_spacing=1.4)

add_rect(s, Inches(0.5), Inches(5.95), Inches(12.33), Inches(0.95), GRAY_BG)
add_rect(s, Inches(0.5), Inches(5.95), Inches(0.1), Inches(0.95), ACCENT)
add_text(s, Inches(0.8), Inches(6.0), Inches(11.8), Inches(0.45),
         "단일 진입점에서 결산 사이클의 모든 단계를 직관적으로 추적 · 관리",
         size=15, bold=True, color=NAVY_DARK)
add_text(s, Inches(0.8), Inches(6.45), Inches(11.8), Inches(0.4),
         "5단계의 결산 워크플로우 + 검증·자동화 라인을 하나의 시스템으로 일원화하여, 어디서 무엇이 막혔는지 즉시 파악 가능",
         size=11, color=GRAY_MID)

# ============================================================
# 슬라이드 7 : 전체 기능 맵 (6개 그룹)
# ============================================================
s = prs.slides.add_slide(blank)
add_title_bar(s, "전체 기능 맵", "6개 기능 그룹 한눈에 보기", 7, TOTAL)

groups = [
    ("A", "메인 / 홈",             ["통합 대시보드", "결산 연도 선택", "진행상태 모니터링"]),
    ("B", "로그인 / 권한 관리",     ["사용자 인증", "관리자 콘솔", "권한 그룹 프리셋"]),
    ("C", "패키지 배포 · 합산",     ["빈 패키지 배포", "일괄 업로드", "자동 합산 + FX"]),
    ("D", "WCE 입력 · 집계",       ["내부거래 입력", "전사 집계", "법인별 관리"]),
    ("E", "연결처리 · 현금정산표", ["연결조정 분개", "다단계 합산", "현금정산표 · KPI"]),
    ("F", "검증 · 자동화",         ["패키지 자동 검증", "COA 누락 감사", "4Q 주석 합산"]),
]
# 3 x 2 그리드
gcw, gch = Inches(4.0), Inches(2.45)
gap_x = Inches(0.15)
gap_y = Inches(0.2)
total_grid_w = gcw*3 + gap_x*2
sx_g = (SW - total_grid_w) / 2
sy_g = Inches(1.75)
colors_g = [NAVY, RGBColor(0x37, 0x4E, 0x8C), RGBColor(0x52, 0x71, 0xAB),
            RGBColor(0x6E, 0x8C, 0xC4), RGBColor(0x89, 0xA6, 0xD8), RGBColor(0xB4, 0x53, 0x09)]
for i, ((tag, title, items), c) in enumerate(zip(groups, colors_g)):
    col, row = i % 3, i // 3
    x = sx_g + col*(gcw + gap_x)
    y = sy_g + row*(gch + gap_y)
    add_rect(s, x, y, gcw, gch, WHITE, line=GRAY_LT)
    add_rect(s, x, y, gcw, Inches(0.75), c)
    add_text(s, x + Inches(0.25), y, Inches(0.7), Inches(0.75),
             tag, size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    add_text(s, x + Inches(0.95), y, gcw - Inches(1.05), Inches(0.75),
             title, size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    for j, it in enumerate(items):
        iy = y + Inches(0.9 + j*0.5)
        add_rect(s, x + Inches(0.3), iy + Inches(0.13), Inches(0.08), Inches(0.18), c)
        add_text(s, x + Inches(0.5), iy + Inches(0.05), gcw - Inches(0.7), Inches(0.4),
                 it, size=11, color=GRAY_DK, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# 기능 소개 슬라이드 공통 함수
# ============================================================
def feature_slide(slide_num, group_tag, group_title, page_title, subtitle,
                  features, business_value, shot_label, shot_keys, new_badge=False):
    s = prs.slides.add_slide(blank)
    add_title_bar(s, page_title, subtitle, slide_num, TOTAL)
    add_rect(s, Inches(0.5), Inches(1.45), Inches(0.55), Inches(0.32), NAVY)
    add_text(s, Inches(0.5), Inches(1.45), Inches(0.55), Inches(0.32),
             group_tag, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.15), Inches(1.45), Inches(7), Inches(0.32),
             group_title, size=12, bold=True, color=GRAY_MID, anchor=MSO_ANCHOR.MIDDLE)
    # NEW 뱃지
    if new_badge:
        add_rect(s, Inches(7.4), Inches(1.45), Inches(0.85), Inches(0.32), ACCENT)
        add_text(s, Inches(7.4), Inches(1.45), Inches(0.85), Inches(0.32),
                 "NEW", size=11, bold=True, color=NAVY_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_screenshot_area(s, Inches(0.5), Inches(2.0), Inches(7.8), Inches(4.7),
                        slide_num, shot_label, keys=shot_keys)

    rx = Inches(8.5)
    rw = Inches(4.35)
    add_text(s, rx, Inches(2.0), rw, Inches(0.4),
             "주요 기능", size=14, bold=True, color=NAVY_DARK)
    add_rect(s, rx, Inches(2.42), Inches(0.4), Emu(28575), ACCENT)
    add_bullets(s, rx, Inches(2.55), rw, Inches(2.6), features, size=12,
                color=GRAY_DK, bullet_color=NAVY, line_spacing=1.35)

    add_rect(s, rx, Inches(5.25), rw, Inches(1.55), NAVY)
    add_text(s, rx + Inches(0.2), Inches(5.35), rw - Inches(0.4), Inches(0.35),
             "도입 효과", size=11, bold=True, color=ACCENT)
    add_text(s, rx + Inches(0.2), Inches(5.65), rw - Inches(0.4), Inches(1.1),
             business_value, size=12, color=WHITE, line_spacing=1.45)
    return s

# ============================================================
# 슬라이드 8 : [A] 메인 / 홈
# ============================================================
feature_slide(
    8, "A", "메인 / 홈",
    "메인 대시보드", "업로드 → 합산 → 경영 대시보드를 한 화면에서 (index.html)",
    [
        "결산 연도 선택 · 일괄 업로드 · 제출 현황 모니터링",
        "합산 결과 미리보기 (회사별 컬럼 + 총합 + 전기 비교)",
        "경영 대시보드 3개 섹션 — 수익성 / 재무안전성 / 현금흐름",
        "회사별 개별 차트 + 그룹 합산 차트 (차입금 구성 · 이익률 비교)",
    ],
    "결산 사이클부터 경영 KPI 시각화까지 하나의 화면에서 처리. 합산 직후 임원진이 별도 자료 없이 즉시 경영 현황을 파악할 수 있습니다.",
    "메인 대시보드 화면", ["main", "home", "index"]
)

# ============================================================
# 슬라이드 9 : [B-1] 로그인 / 비밀번호
# ============================================================
feature_slide(
    9, "B", "로그인 / 권한 관리",
    "로그인 및 비밀번호 정책", "안전한 사용자 인증 · 비밀번호 변경 (login.html)",
    [
        "사용자 ID · 비밀번호 기반 인증",
        "비밀번호 해시 저장 (Scrypt 알고리즘)",
        "사용자별 비밀번호 자율 변경 화면 제공",
        "관리자에 의한 비밀번호 초기화 지원",
    ],
    "비밀번호 평문 저장을 차단하고, 회사 보안 정책과 정합성 있는 인증 체계를 제공합니다.",
    "로그인 / 비밀번호 변경 화면", ["login", "로그인"]
)

# ============================================================
# 슬라이드 10 : [B-2] 사용자 · 권한 관리
# ============================================================
feature_slide(
    10, "B", "로그인 / 권한 관리",
    "사용자 · 권한 관리", "관리자 콘솔과 권한 그룹 프리셋 (admin_users.html / admin_permission_groups.html)",
    [
        "사용자 생성 · 삭제 · 비밀번호 초기화",
        "자회사 단위의 접근 권한 부여 · 회수",
        "권한 그룹(프리셋) 기반 일괄 적용",
        "신규 담당자 온보딩 절차 표준화",
    ],
    "자회사별 데이터 분리로 정보 접근 통제를 강화하고, 권한 그룹 프리셋으로 신규 담당자 온보딩 시간을 단축합니다.",
    "사용자 · 권한 관리 화면", ["admin_user", "admin_permission", "권한", "permission"]
)

# ============================================================
# 슬라이드 11 : [C-1] 패키지 업로드
# ============================================================
feature_slide(
    11, "C", "패키지 배포 · 합산",
    "결산패키지 업로드", "자회사 결산자료 일괄 수집 (index.html — 업로드 영역)",
    [
        "다수의 .xlsx / .xlsm 패키지 동시 업로드",
        "결산 연도 · 기간별 자동 보관",
        "업로드 즉시 양식 적합성 1차 검증",
        "자회사별 제출 현황 시각화",
    ],
    "이메일로 분산 수집되던 패키지를 시스템 단일 채널로 통합하여, 누락 · 중복 · 양식 오류를 사전에 차단합니다.",
    "패키지 업로드 화면", ["upload", "업로드", "package"]
)

# ============================================================
# 슬라이드 12 : [C-2] 자동 합산 결과
# ============================================================
feature_slide(
    12, "C", "패키지 배포 · 합산",
    "자동 합산", "BS · PL · 현금흐름표 자동 합산 (합산 실행 · 결과 화면)",
    [
        "BS · PL · 현금흐름표 동시 합산",
        "자산 · 부채 · 자본 · 차입금 · 부채비율 등 핵심지표 미리보기",
        "WCE(내부거래 제거) 결과 자동 반영",
        "통합합산결과 + 회사별 전기비교 엑셀 자동 생성 · 다운로드",
    ],
    "분기당 수일 단위의 수기 합산 공수를 제거하고, 합산 단계에서 발생하던 휴먼 에러를 구조적으로 차단합니다.",
    "자동 합산 결과 화면", ["aggregate", "합산", "consolidate"]
)

# ============================================================
# 슬라이드 13 : [C-3] 환율 적용
# ============================================================
feature_slide(
    13, "C", "패키지 배포 · 합산",
    "환율(FX) 일괄 적용", "통화 환산 자동화 및 재적용 (환율 관리 영역)",
    [
        "분기별 · 통화별 환율 일괄 등록",
        "BS(기말환율 spot) · PL(평균환율 avg) 자동 구분 적용",
        "전기 분기말 환율 자동 제안 (입력 편의)",
        "환율 변경 시 1-Click 일괄 재합산 (reapply)",
    ],
    "통화 환산 오류와 재작업 비용을 차단합니다. 환율 변경이 발생해도 즉시 일괄 재계산되어 결산 적시성을 확보합니다.",
    "환율 관리 화면", ["fx", "환율", "rate"]
)

# ============================================================
# 슬라이드 14 : [C-4] 분기 패키지 자동 배포  [NEW]
# ============================================================
feature_slide(
    14, "C", "패키지 배포 · 합산",
    "분기 패키지 자동 배포", "빈 결산패키지 일괄 생성 · 자가 다운로드 (admin_distribute.html)",
    [
        "연도별 빈 .xlsm 템플릿 · 분기별 시트보호 비밀번호 등록",
        "전년 동기 BS / PL 자동 자기참조 (이월값 자동 채움)",
        "중앙환율 자동 적용 · WCE 자본 잔액 자동 산정",
        "자회사 담당자 본인 패키지 자가 다운로드",
    ],
    "분기 시작 시 자회사에 빈 패키지를 일관된 비밀번호 정책으로 일괄 배포하고, 양식 · 환율 · 전기 잔액 누락 위험을 차단합니다.",
    "패키지 배포 관리 화면", ["distribute", "배포", "패키지배포"],
    new_badge=True
)

# ============================================================
# 슬라이드 15 : [D-1] WCE 입력
# ============================================================
feature_slide(
    15, "D", "WCE 입력 · 집계",
    "WCE 입력", "내부거래 제거 데이터 입력 (wce_input.html · wce_list.html)",
    [
        "자회사별 WCE(Within-Company Elimination) 입력",
        "5개 자본 항목 표 (자본금 · 이익잉여금 등) 별도 입력",
        "연도 · 법인별 분리 저장 및 회사별 목록 조회",
        "전기 Q4 기말 → 당기 기초 자동 산정 (비-첫해)",
    ],
    "자회사 간 내부거래 데이터를 표준화된 양식으로 수집하여, 연결조정의 일관성을 확보합니다.",
    "WCE 입력 화면", ["wce_input", "wce입력", "wce_list", "wce목록"]
)

# ============================================================
# 슬라이드 16 : [D-2] WCE 통합 집계
# ============================================================
feature_slide(
    16, "D", "WCE 입력 · 집계",
    "WCE 통합 집계", "전사 단위 내부거래 집계 (wce_aggregate.html)",
    [
        "전 자회사 WCE 데이터 통합 집계 뷰",
        "자본 항목 표별 · 회사별 · 계정 코드별 매트릭스",
        "Local 통화와 KRW 환산금액 동시 표시",
        "환산효과 제외 기말금액 자동 계산 (ending)",
    ],
    "전사 단위에서 내부거래 흐름을 즉시 파악할 수 있어, 연결조정 단계의 분석 · 의사결정 속도를 높입니다.",
    "WCE 통합 집계 화면", ["wce_aggregate", "wce집계"]
)

# ============================================================
# 슬라이드 17 : [E-1] 연결 그룹 · 다단계 구조
# ============================================================
feature_slide(
    17, "E", "연결처리 · 현금정산표",
    "연결 그룹 · 다단계 구조", "복잡한 지배구조의 체계적 표현 (consolidation.html)",
    [
        "연결 그룹 정의 · 관리 (consol_groups.json)",
        "다단계 지배구조 (예: 글로벌세아 ⊃ 상역 ⊃ 태림/GIT)",
        "그룹정보.xlsx를 통한 회사 · 그룹 일괄 등록",
        "전년 동기 값 입력으로 YoY 비교 기반 마련",
    ],
    "다단계 지배구조를 시스템에 명시적으로 표현함으로써, 합산 누락 · 중복 위험을 구조적으로 제거합니다.",
    "연결 그룹 관리 화면", ["consolidation", "연결그룹", "group"]
)

# ============================================================
# 슬라이드 18 : [E-2] 분개 자동 검증
# ============================================================
feature_slide(
    18, "E", "연결처리 · 현금정산표",
    "연결조정 분개 자동 검증", "차변 · 대변 균형, 코드 누락 차단 (consolidation.html)",
    [
        "연결조정 분개 엑셀 일괄 업로드",
        "차변 · 대변 균형 자동 검증",
        "계정 코드 누락 · 미정의 코드 차단",
        "그룹 재귀 합산으로 누락 · 중복 제거",
    ],
    "분개 검증 자동화로 사후 정정 비용을 절감하고, 결산 신뢰성을 결정적으로 끌어올립니다.",
    "분개 검증 화면", ["journal", "분개", "consol_verify"]
)

# ============================================================
# 슬라이드 19 : [E-3] 연결 KPI 대시보드
# ============================================================
feature_slide(
    19, "E", "연결처리 · 현금정산표",
    "연결 KPI 대시보드", "임원용 한눈에 보는 결산 결과 (consolidation_dashboard.html)",
    [
        "최종 연결재무제표 엑셀 자동 생성",
        "그룹별 매출 · 이익 · 자산 KPI 시각화",
        "도넛 · 막대 차트 기반 구성 비중 표시",
        "전년 동기 대비 YoY 추이 자동 비교",
    ],
    "결산 직후 즉시 시각화되어 경영진의 적시 의사결정을 지원합니다. 별도 보고자료 작성 공수 없이 화면 그대로 활용 가능합니다.",
    "연결 KPI 대시보드", ["dashboard", "대시보드", "kpi"]
)

# ============================================================
# 슬라이드 20 : [E-4] 현금정산표  [NEW]
# ============================================================
feature_slide(
    20, "E", "연결처리 · 현금정산표",
    "현금정산표", "연결정산표 → 현금흐름표 1-Click 산출 (cash_worksheet.html)",
    [
        "연결그룹별 BS · PL · CF 합산 + 연결조정 · 내부거래 분개 자동 반영",
        "수기조정 · 단수조정 · 코멘트 인라인 입력 (JSON 영속화)",
        "NI plug 보정으로 연결정산표 ↔ CF의 NI 일치 자동 보장",
        "항등식 검증 뱃지 · 그룹 자금조정 별도 컬럼 · 엑셀 다운로드",
    ],
    "연결정산표에서 현금흐름표까지 한 번의 클릭으로 산출하고, NI 정합을 자동 검증하여 분개 정정 비용을 절감합니다.",
    "현금정산표 화면", ["cash_worksheet", "현금정산", "cashflow"],
    new_badge=True
)

# ============================================================
# 슬라이드 21 : [E-5] CF 매핑 편집기  [NEW]
# ============================================================
feature_slide(
    21, "E", "연결처리 · 현금정산표",
    "CF 매핑 편집기", "계정코드 → 현금흐름 라인 매핑 일괄 편집 (cash_worksheet_mapping.html)",
    [
        "cf_mapping_v2_draft.json의 COA 행별 매핑을 화면에서 편집",
        "연결조정 · 내부거래 각각 CF 코드 · 부호 토글 인터페이스",
        "cf_lines 자동완성 · 변경분 시각화 · 정책 잠금 표시",
        "FileLock 안전 저장 · .bak 자동 백업으로 사고 방지",
    ],
    "매핑을 JSON 직접 편집할 필요 없이 화면에서 검토 · 수정 가능. 신규 계정 추가 시 매핑 정비 비용과 실수 가능성을 줄입니다.",
    "CF 매핑 편집기", ["cf_mapping", "매핑편집", "mapping"],
    new_badge=True
)

# ============================================================
# 슬라이드 22 : [F-1] 패키지 자동 검증  [NEW]
# ============================================================
feature_slide(
    22, "F", "검증 · 자동화",
    "패키지 자동 검증", "업로드된 WCF 데이터 4종 자동 검사 (admin/package-verify)",
    [
        "WCF Diff 검증 : PL · MF 불일치 행 + 원인 미기재 검출",
        "계정 부호 검증 : 비용 칸 수익 계정 · 수익 칸 비용 계정 검출",
        "흐름 부호 검증 : Cash in/out 합계 부호 자동 점검",
        "퇴직금 부호 검증 (CF2200401 별도 정책 적용)",
    ],
    "업로드 직후 패키지 데이터의 정합성을 사전 점검하여, 연결조정 단계에서 오류를 추적하는 비용을 사전에 차단합니다.",
    "패키지 자동 검증 화면", ["package_verify", "패키지검증", "verify"],
    new_badge=True
)

# ============================================================
# 슬라이드 23 : [F-2] COA 누락 코드 검출  [NEW]
# ============================================================
feature_slide(
    23, "F", "검증 · 자동화",
    "COA 누락 코드 검출", "신규 계정 코드 자동 식별 · 일괄 등록 (admin/coa-audit)",
    [
        "BS · PL_MF · CF의 모든 ref_code 스캔 → 미등록 코드 검출",
        "prefix 기반 추천 섹션 자동 산정 (Ⅰ-2 · Ⅰ-3 등)",
        "‘NI 흡수’ · ‘자체 라인 권장’ · ‘사례별’ 가이드 표시",
        "선택분만 cf_lines 일괄 추가 · prefix 일관성 강제 · .bak 자동 백업",
    ],
    "신규 계정 누락으로 인한 합산 불일치를 사전에 차단하고, 코드 등록 작업을 클릭 한 번으로 처리할 수 있습니다.",
    "COA Audit 화면", ["coa_audit", "coa", "코드감사"],
    new_badge=True
)

# ============================================================
# 슬라이드 24 : [F-3] 4Q 주석 합산  [NEW]
# ============================================================
feature_slide(
    24, "F", "검증 · 자동화",
    "4Q 주석 합산", "단기차입금 등 주석 자료 자동 집계 (admin/note-aggregate)",
    [
        "4Q 분기 패키지 L1 시트 일괄 스캔",
        "회사별 currency · spot 환율 적용해 KRW 환산",
        "종류별 · 회사별 매트릭스 자동 재배열",
        "결과 엑셀 자동 생성 · 다운로드 링크 제공",
    ],
    "4Q 주석 작성을 위한 차입금 정보의 수동 취합을 완전 자동화하여, 결산 마감 단계의 공수를 추가로 절감합니다.",
    "4Q 주석 합산 화면", ["note_aggregate", "주석", "l1_borrowing"],
    new_badge=True
)

# ============================================================
# 슬라이드 25 : 정량적 기대효과
# ============================================================
s = prs.slides.add_slide(blank)
add_title_bar(s, "정량적 기대효과", "도입 전 → 도입 후 비교", 25, TOTAL)

kpis = [
    ("분기당 합산 공수", "수일~수주",       "수십 분",         "▼ 약 90%+", NAVY),
    ("환율 변경 재작업", "전체 수기 재작업", "1-Click 재합산",  "즉시 처리",  RGBColor(0x0D, 0x94, 0x88)),
    ("분개 균형 검증",   "수기 점검",        "자동 100% 검증",  "오류율 ↓",   RGBColor(0xB4, 0x53, 0x09)),
    ("결산 결과 시각화", "별도 자료 작성",   "실시간 대시보드", "보고 공수 ↓", RGBColor(0x8B, 0x5C, 0xF6)),
]
cw, ch = Inches(2.95), Inches(4.4)
gap = Inches(0.15)
total_w = cw*4 + gap*3
sx = (SW - total_w) / 2
sy = Inches(1.8)
for i, (title, before, after, delta, c) in enumerate(kpis):
    x = sx + i*(cw + gap)
    add_rect(s, x, sy, cw, ch, WHITE, line=GRAY_LT)
    add_rect(s, x, sy, cw, Inches(0.7), c)
    add_text(s, x, sy + Inches(0.12), cw, Inches(0.5),
             title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), sy + Inches(0.95), cw - Inches(0.4), Inches(0.3),
             "AS-IS", size=10, bold=True, color=GRAY_MID, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), sy + Inches(1.25), cw - Inches(0.4), Inches(0.5),
             before, size=14, color=GRAY_DK, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, x, sy + Inches(1.95), cw, Inches(0.4),
             "▼", size=18, color=c, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, x + Inches(0.2), sy + Inches(2.4), cw - Inches(0.4), Inches(0.3),
             "TO-BE", size=10, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), sy + Inches(2.7), cw - Inches(0.4), Inches(0.5),
             after, size=14, color=NAVY_DARK, align=PP_ALIGN.CENTER, bold=True)
    add_rect(s, x + Inches(0.3), sy + Inches(3.55), cw - Inches(0.6), Inches(0.5), c)
    add_text(s, x + Inches(0.3), sy + Inches(3.55), cw - Inches(0.6), Inches(0.5),
             delta, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
         "※ 효과 수치는 도입 단계에서 사용 환경 · 데이터 양에 따라 조정됩니다.",
         size=10, color=GRAY_MID, italic=True, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 26 : 정성적 기대효과
# ============================================================
s = prs.slides.add_slide(blank)
add_title_bar(s, "정성적 기대효과", "수치로 측정하기 어려운 구조적 가치", 26, TOTAL)

qual = [
    ("결산 신뢰성 향상",    "균형 검증 · 코드 검증 · 패키지 검증 자동화로 결산 품질 표준화"),
    ("절차 표준화",         "동일 양식 · 동일 계산식 · 동일 배포 정책으로 회사 간 정합성 확보"),
    ("노하우 자산화",       "결산 노하우(CF 매핑·검증 룰 등)를 시스템에 내재화하여 속인화 해소"),
    ("감사 대응력 강화",    "권한 분리 · 이력 관리 · 표준 절차로 감사 응답 비용 감소"),
    ("의사결정 속도 향상",  "결산 직후 즉시 KPI 시각화 → 적시 경영판단 지원"),
    ("협업 환경 개선",      "단일 시스템에서 자회사 · 본사 · 감사가 동일 데이터 참조"),
]
cw, ch = Inches(6.0), Inches(1.4)
gap_x, gap_y = Inches(0.3), Inches(0.2)
sx, sy = Inches(0.55), Inches(1.7)
for i, (t, d) in enumerate(qual):
    col, row = i % 2, i // 2
    x = sx + col*(cw + gap_x)
    y = sy + row*(ch + gap_y)
    add_rect(s, x, y, cw, ch, WHITE, line=GRAY_LT)
    add_rect(s, x, y, Inches(0.08), ch, NAVY)
    add_rect(s, x + Inches(0.3), y + Inches(0.4), Inches(0.6), Inches(0.6), ICE)
    add_text(s, x + Inches(0.3), y + Inches(0.4), Inches(0.6), Inches(0.6),
             "✓", size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.05), y + Inches(0.2), cw - Inches(1.2), Inches(0.45),
             t, size=15, bold=True, color=NAVY_DARK)
    add_text(s, x + Inches(1.05), y + Inches(0.72), cw - Inches(1.2), Inches(0.6),
             d, size=12, color=GRAY_MID)

# ============================================================
# 슬라이드 27 : 기술 스택 및 보안
# ============================================================
s = prs.slides.add_slide(blank)
add_title_bar(s, "기술 스택 및 보안", "내부 폐쇄망 운영을 전제로 한 설계", 27, TOTAL)

add_rect(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(5.3), WHITE, line=GRAY_LT)
add_rect(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(0.55), NAVY)
add_text(s, Inches(0.7), Inches(1.6), Inches(5.9), Inches(0.55),
         "기술 스택", size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
tech = [
    ("백엔드",       "Python 3 / Flask"),
    ("데이터 처리",  "openpyxl · pandas 기반 엑셀 자동화"),
    ("프론트엔드",   "HTML · CSS · JavaScript (서버 렌더링)"),
    ("저장 형식",    "JSON 파일 기반 메타데이터 + Excel 산출물"),
    ("실행 환경",    "Windows 서버, 표준 .bat 실행 스크립트"),
    ("외부 의존성",  "외부 인터넷 통신 없음 (폐쇄망)"),
]
for i, (k, v) in enumerate(tech):
    y = Inches(2.35 + i*0.7)
    add_text(s, Inches(0.7), y, Inches(1.7), Inches(0.4),
             k, size=12, bold=True, color=NAVY_DARK)
    add_text(s, Inches(2.4), y, Inches(4.1), Inches(0.4),
             v, size=12, color=GRAY_DK)
    if i < len(tech)-1:
        add_rect(s, Inches(0.7), y + Inches(0.5), Inches(5.7), Emu(9525), GRAY_LT)

add_rect(s, Inches(6.85), Inches(1.6), Inches(5.95), Inches(5.3), WHITE, line=GRAY_LT)
add_rect(s, Inches(6.85), Inches(1.6), Inches(5.95), Inches(0.55), NAVY_DARK)
add_text(s, Inches(7.05), Inches(1.6), Inches(5.75), Inches(0.55),
         "보안 설계", size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
sec = [
    "비밀번호 해시 저장 (평문 미저장)",
    "자회사별 데이터 접근 권한 분리",
    "사용자 · 권한 · 합산 이력 로그 보관",
    "외부 인터넷 통신 차단 (폐쇄망 전제)",
    "표준 .bat 실행으로 권한 통제 단순화",
    "감사 대응을 위한 변경 이력 추적",
]
for i, item in enumerate(sec):
    y = Inches(2.5 + i*0.7)
    add_rect(s, Inches(7.1), y + Inches(0.15), Inches(0.18), Inches(0.18), ACCENT)
    add_text(s, Inches(7.4), y, Inches(5.3), Inches(0.5),
             item, size=12, color=GRAY_DK, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# 슬라이드 28 : 운영 사양 및 배포
# ============================================================
s = prs.slides.add_slide(blank)
add_title_bar(s, "운영 사양 및 배포", "서버 요구 사항과 배포 방식", 28, TOTAL)

add_text(s, Inches(0.55), Inches(1.7), Inches(6), Inches(0.5),
         "서버 요구 사양 (최소 / 권장)", size=15, bold=True, color=NAVY_DARK)
specs = [
    ["구분",     "최소 사양",        "권장 사양"],
    ["OS",       "Windows Server",   "Windows Server 2019+"],
    ["CPU",      "4 Core",           "8 Core"],
    ["RAM",      "8 GB",             "16 GB"],
    ["저장공간", "100 GB",           "500 GB (이력 누적 대비)"],
    ["네트워크", "사내 폐쇄망",      "사내 폐쇄망 (외부 차단)"],
]
table_x = Inches(0.55)
table_y = Inches(2.3)
col_ws = [Inches(1.6), Inches(2.2), Inches(2.4)]
row_h = Inches(0.5)
for ri, row in enumerate(specs):
    is_head = (ri == 0)
    for ci, cell in enumerate(row):
        x = table_x + sum(col_ws[:ci], Emu(0))
        y = table_y + ri * row_h
        if is_head:
            add_rect(s, x, y, col_ws[ci], row_h, NAVY)
            add_text(s, x, y, col_ws[ci], row_h, cell, size=12, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            bg = WHITE if ri % 2 == 1 else GRAY_BG
            add_rect(s, x, y, col_ws[ci], row_h, bg, line=GRAY_LT)
            bold_c = (ci == 0)
            add_text(s, x + Inches(0.15), y, col_ws[ci] - Inches(0.3), row_h, cell,
                     size=12, color=GRAY_DK, bold=bold_c, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(7.4), Inches(1.7), Inches(5.4), Inches(5.0), WHITE, line=GRAY_LT)
add_rect(s, Inches(7.4), Inches(1.7), Inches(0.1), Inches(5.0), NAVY)
add_text(s, Inches(7.65), Inches(1.85), Inches(5), Inches(0.5),
         "배포 · 운영 방식", size=15, bold=True, color=NAVY_DARK)
add_rect(s, Inches(7.65), Inches(2.4), Inches(0.5), Emu(28575), ACCENT)

deploy = [
    ("단순 배포",      "Python 환경 + 폴더 복사로 신규 서버 즉시 운영 가능"),
    ("1-Click 실행",   "setup_and_run.bat 더블클릭으로 서비스 기동"),
    ("브라우저 접속",  "사내 PC에서 표준 웹브라우저로 즉시 사용"),
    ("백업 단순화",    "전체 폴더 단위 백업 / 복원 (DB 별도 없음)"),
    ("운영 부담 최소", "전담 운영자 없이도 안정적 운영 가능"),
]
for i, (t, d) in enumerate(deploy):
    y = Inches(2.7 + i*0.85)
    add_text(s, Inches(7.65), y, Inches(5), Inches(0.4),
             t, size=13, bold=True, color=NAVY_DARK)
    add_text(s, Inches(7.65), y + Inches(0.42), Inches(5), Inches(0.4),
             d, size=11, color=GRAY_MID)

# ============================================================
# 슬라이드 29 : 협조 요청 사항
# ============================================================
s = prs.slides.add_slide(blank)
add_title_bar(s, "IT부서 협조 요청 사항", "도입 · 운영을 위한 협의 항목", 29, TOTAL)

reqs = [
    ("01", "운영 서버 할당",     "사내 폐쇄망에 Windows 서버 1대 할당 및 접근 권한 설정"),
    ("02", "방화벽 · DNS 설정",  "사내 DNS 등록 또는 IP 기반 접속 경로 설정"),
    ("03", "사용자 계정 정책",   "회사 계정 정책과 비밀번호 정책 정합성 협의"),
    ("04", "백업 정책",          "전체 폴더 단위 정기 백업 일정 및 보관 기간 협의"),
    ("05", "보안 점검 협조",     "최초 도입 시 보안 검토 및 사용자 권한 가이드 수립"),
    ("06", "장기 운영 인수인계", "전담 IT 운영자 지정 및 운영 매뉴얼 인수인계"),
]
sx = Inches(0.5)
sy = Inches(1.65)
cw = Inches(12.33)
ch = Inches(0.78)
gap = Inches(0.12)
for i, (n, t, d) in enumerate(reqs):
    y = sy + i*(ch + gap)
    add_rect(s, sx, y, cw, ch, WHITE, line=GRAY_LT)
    add_rect(s, sx, y, Inches(0.95), ch, NAVY)
    add_text(s, sx, y, Inches(0.95), ch,
             n, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    add_text(s, sx + Inches(1.2), y + Inches(0.1), Inches(3.5), Inches(0.55),
             t, size=14, bold=True, color=NAVY_DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, sx + Inches(4.9), y + Inches(0.1), cw - Inches(5.2), Inches(0.55),
             d, size=12, color=GRAY_DK, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# 슬라이드 30 : 도입 로드맵
# ============================================================
s = prs.slides.add_slide(blank)
add_title_bar(s, "도입 로드맵", "단계별 일정", 30, TOTAL)

phases = [
    ("Phase 1", "준비",          "1주",   "서버 할당 · 환경 구축 · 초기 데이터 이관"),
    ("Phase 2", "파일럿 운영",   "2~3주", "일부 자회사 대상 시범 운영 · 피드백 반영"),
    ("Phase 3", "전사 확대",     "1개월", "전 자회사 대상 단계적 확대 · 사용자 교육"),
    ("Phase 4", "정식 운영",     "지속",   "분기 결산 정식 적용 · 지속 개선"),
]
add_rect(s, Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.04), ICE)

cw = Inches(2.7)
gap = Inches(0.1)
total_w = cw*4 + gap*3
sx = (SW - total_w) / 2

for i, (ph, title, dur, desc) in enumerate(phases):
    x = sx + i*(cw + gap)
    add_rect(s, x, Inches(1.8), cw, Inches(1.5), WHITE, line=GRAY_LT)
    add_rect(s, x, Inches(1.8), cw, Inches(0.4), NAVY)
    add_text(s, x, Inches(1.8), cw, Inches(0.4),
             ph, size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, Inches(2.3), cw, Inches(0.5),
             title, size=18, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.85), cw, Inches(0.35),
             f"⏱ {dur}", size=11, color=GRAY_MID, align=PP_ALIGN.CENTER, italic=True)
    mx = x + cw/2 - Inches(0.15)
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, mx, Inches(3.35), Inches(0.3), Inches(0.3))
    shp.fill.solid(); shp.fill.fore_color.rgb = ACCENT
    shp.line.color.rgb = ACCENT
    add_rect(s, x, Inches(3.95), cw, Inches(2.0), GRAY_BG)
    add_text(s, x + Inches(0.2), Inches(4.1), cw - Inches(0.4), Inches(1.8),
             desc, size=12, color=GRAY_DK, line_spacing=1.45, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.75), NAVY)
add_text(s, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.75),
         "약 2개월 내 전사 정식 운영 진입을 목표로 합니다.",
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# 슬라이드 31 : Q&A / 마무리
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SW, SH, NAVY_DARK)
add_rect(s, 0, 0, Inches(0.18), SH, ACCENT)
add_text(s, 0, Inches(2.1), SW, Inches(1.8),
         "Q & A", size=100, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_HEAD)
add_rect(s, Inches(6.2), Inches(4.4), Inches(0.95), Emu(38100), ACCENT)
add_text(s, 0, Inches(4.65), SW, Inches(0.5),
         "Thank You", size=22, color=ICE, align=PP_ALIGN.CENTER, italic=True)
add_text(s, 0, Inches(5.4), SW, Inches(0.5),
         "연결재무보고 통합 시스템 | 임원 보고용", size=13, color=ICE, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(6.0), SW, Inches(0.5),
         "재무팀  |  2026", size=11, color=ICE, align=PP_ALIGN.CENTER)

# ============================================================
prs.save(OUTPUT_PATH)
print(f"[OK] PPT generated: {OUTPUT_PATH}")
print(f"[OK] Slide count: {len(prs.slides)}")
print(f"[OK] Screenshots dir: {SHOTS_DIR}")
